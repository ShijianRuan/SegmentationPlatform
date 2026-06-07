#!/usr/bin/env python3
"""Generate Segmentation Platform Architecture PPTX using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
DARK  = RGBColor(0x0F, 0x17, 0x2A)
BLUE  = RGBColor(0x02, 0x84, 0xC7)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED   = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF1, 0xF5, 0xF9)
GRAY  = RGBColor(0x94, 0xA3, 0xB8)
DGRAY = RGBColor(0x47, 0x55, 0x69)
BLACK = RGBColor(0x1E, 0x29, 0x3B)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ──
def add_blank_slide():
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color=LGRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Arial', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=BLACK, font_name='Arial', line_spacing=1.3):
    """lines: list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, fs, clr = item, False, font_size, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            fs = item[2] if len(item) > 2 else font_size
            clr = item[3] if len(item) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(fs)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.font.name = font_name
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        try:
            p.line_spacing = Pt(int(fs * line_spacing))
        except:
            pass
    return txBox

def slide_title_bar(slide, title_text, subtitle_text=None):
    """Dark title bar at top"""
    bar = add_rect(slide, 0, 0, W, Inches(1.15), fill_color=DARK)
    add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.5), title_text,
             font_size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle_text:
        add_text(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.4), subtitle_text,
                 font_size=13, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
    # Accent line
    add_rect(slide, 0, Inches(1.15), W, Inches(0.04), fill_color=BLUE)

def add_table(slide, left, top, col_widths, rows, header_bg=DARK, has_header=True):
    """rows: list of lists. Each cell: str or (text, bold, color, bg_color)"""
    n_rows = len(rows)
    n_cols = len(rows[0])
    row_h = Inches(0.42)
    table_width = sum(col_widths)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, row_h * n_rows)
    tbl = tbl_shape.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = row_h
        for ci, cell_val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Parse cell value
            if isinstance(cell_val, str):
                txt, bld, clr, bg = cell_val, False, BLACK, None
            else:
                txt = cell_val[0]
                bld = cell_val[1] if len(cell_val) > 1 else False
                clr = cell_val[2] if len(cell_val) > 2 else BLACK
                bg = cell_val[3] if len(cell_val) > 3 else None
            # Cell fill
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            elif has_header and ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif ri % 2 == 0 and ri > 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LGRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            # Cell text
            p = cell.text_frame.paragraphs[0]
            p.text = txt
            p.font.size = Pt(10)
            p.font.bold = bld or (has_header and ri == 0)
            p.font.color.rgb = WHITE if (has_header and ri == 0) else clr
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.LEFT
            try:
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
            except:
                pass
    return tbl_shape

# ═══════════════════════════════════════
# S1: Title Slide
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, DARK)
add_rect(s, 0, 0, W, Inches(0.06), fill_color=BLUE)
add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0), "医学图像分割平台架构设计",
         font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_rect(s, Inches(4.5), Inches(3.1), Inches(4.3), Inches(0.04), fill_color=BLUE)
add_text(s, Inches(1), Inches(3.4), Inches(11), Inches(0.6), "Segmentation Platform Architecture Design",
         font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4), "2026-06-07",
         font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# S2: Core Positioning
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "平台核心定位")

# Cards in a row
cards = [
    ("病例\nCase", BLUE),
    ("图像\nImage Artifact", BLUE),
    ("标签\nLabel Artifact", GREEN),
    ("训练任务\nTask", AMBER),
    ("模型版本\nModel Record", BLUE),
]
cw = Inches(2.2)
ch = Inches(2.6)
start_x = Inches(0.5)
gap = Inches(0.3)
for i, (label, col) in enumerate(cards):
    x = start_x + i * (cw + gap)
    y = Inches(1.8)
    c = add_rect(s, x, y, cw, ch, fill_color=WHITE, line_color=col, line_width=Pt(1.5))
    add_text(s, x + Inches(0.1), y + Inches(0.15), cw - Inches(0.2), ch - Inches(0.3),
             label, font_size=15, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(cards) - 1:
        add_text(s, x + cw - Inches(0.05), y + Inches(1.1), Inches(0.4), Inches(0.3),
                 "→", font_size=24, color=BLUE, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom insight
bx = add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6), fill_color=WHITE)
lines = [
    ("平台中心 ≠ 某个训练框架或标注软件", True, 18, BLUE),
    ("", False, 6, BLACK),
    ("平台中心 = 病例、图像、标签、训练任务、模型版本之间的关系", True, 16, BLACK),
    ("标注、训练、伪标签、离线推理都围绕这些关系发生", False, 13, DGRAY),
]
add_multiline(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.3), lines)

# ═══════════════════════════════════════
# S3: Three Domains
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "三大实现域", "按「不同责任」划分，不是按 pipeline 顺序")

domains = [
    ("labeling  标注域", "标签生产与安全导回", "导出 Case Package\nMimics 修正 + 保存\n几何校验\n注册 verified_label", BLUE),
    ("training  训练域", "任务消费标签→产出模型", "TaskLabelMap\nDataset Snapshot\nnnUNet Adapter\nModel Record", GREEN),
    ("label_generation  标签生成域", "候选标签生成与回流治理", "公开算法 / 自训模型推理\ncandidate_label\nQC + 准入策略\n回流到标注或训练", AMBER),
]
dw = Inches(3.8)
gap = Inches(0.35)
for i, (title, subtitle, items, col) in enumerate(domains):
    x = Inches(0.5) + i * (dw + gap)
    y = Inches(1.6)
    # Card
    add_rect(s, x, y, dw, Inches(5.0), fill_color=WHITE, line_color=col, line_width=Pt(1.5))
    # Accent bar
    add_rect(s, x, y, Inches(0.08), Inches(5.0), fill_color=col)
    # Title
    add_text(s, x + Inches(0.25), y + Inches(0.2), dw - Inches(0.35), Inches(0.4),
             title, font_size=16, bold=True, color=BLACK)
    # Subtitle
    add_text(s, x + Inches(0.25), y + Inches(0.7), dw - Inches(0.35), Inches(0.35),
             subtitle, font_size=12, bold=True, color=DGRAY)
    # Items
    add_text(s, x + Inches(0.25), y + Inches(1.2), dw - Inches(0.35), Inches(3.5),
             items, font_size=11, color=DGRAY)

# ═══════════════════════════════════════
# S4: Closed Loop Flow
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "闭环流程")

# 8 node table-style flow
nodes = [
    ("病例数据", "标签来源\n(人工/模型/算法)", "标注审核\n(labeling)", "数据集快照\n(Dataset Snapshot)"),
    ("候选标签\n(label_generation)", "批量推理", "模型\n(Model Record)", "训练\n(training)"),
]
nw = Inches(2.6)
nh = Inches(1.0)
for row_idx, row_nodes in enumerate(nodes):
    for col_idx, label in enumerate(row_nodes):
        if row_idx == 0:
            x = Inches(0.5) + col_idx * (nw + Inches(0.4))
            y = Inches(1.6)
        else:
            x = Inches(9.8) - col_idx * (nw + Inches(0.4))
            y = Inches(4.6)
        add_rect(s, x, y, nw, nh, fill_color=WHITE, line_color=BLUE, line_width=Pt(1))
        add_text(s, x + Inches(0.1), y + Inches(0.05), nw - Inches(0.2), nh - Inches(0.1),
                 label, font_size=12, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════
# S5: Label States
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "六种标签状态")

# State table
rows = [
    [("状态", True, WHITE), ("含义", True, WHITE), ("默认可训练", True, WHITE)],
    ["source_label", "外部数据集自带标签，未经平台判断", "视来源质量"],
    ["candidate_label", "模型/算法生成候选结果", "默认否，策略可纳入"],
    ["draft_label", "专门准备给人修正的起点", "否"],
    ["accepted_pseudo_label", "经 QC + 策略接受的伪标签", "任务策略决定"],
    ["verified_label", "人工确认标签（单人保存即 verified）", "是"],
    ["rejected_label", "已判定不可用", "否"],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(3.0), Inches(6.5), Inches(2.8)], rows)

# Transition diagram as text
lines = [
    ("标签状态转换", True, 16, BLACK),
    ("source_label → verified_label（可信源直接确认）", False, 13, DGRAY),
    ("candidate_label → draft_label → 人工修正 → verified_label", False, 13, DGRAY),
    ("candidate_label →（QC 通过）→ accepted_pseudo_label", False, 13, DGRAY),
    ("candidate_label →（QC 失败）→ rejected_label", False, 13, DGRAY),
]
add_multiline(s, Inches(0.5), Inches(4.8), Inches(12), Inches(2.0), lines)

# ═══════════════════════════════════════
# S6: 3-Layer Label Map
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "三层 Label Map 设计")

rows = [
    [("层级", True, WHITE), ("回答什么问题", True, WHITE), ("约束来源", True, WHITE), ("liver 的编号示例", True, WHITE)],
    ["anatomy_vocabulary （语义层）", "这是什么器官？", "无约束，纯名称", "liver = 肝脏"],
    ["review_label_map （工具层）", "工具里是几号？", "标注软件限制", "liver = 10（Mimics）"],
    ["task_label_maps （训练层）", "训练时是几号？", "nnUNet 要求连续整数", "liver = 2（CT5_Liver）"],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(3.5), Inches(3.0), Inches(3.0), Inches(3.0)], rows)

# Merge table
rows2 = [
    [("尝试合并方案", True, WHITE), ("为什么不行", True, WHITE)],
    ["合并 anatomy + review", "换标注工具就要改器官名称表（如 Mimics → 3D Slicer）"],
    ["合并 review + task（最诱人）", "同一 liver 在不同任务编号不同 (2 vs 5 vs 37)，没有唯一编号"],
    ["全合并成一张表", "语义层、工具层、训练层的约束来源各不相同，修改互相拉扯"],
]
add_table(s, Inches(0.5), Inches(4.3), [Inches(4.0), Inches(8.5)], rows2)

# ═══════════════════════════════════════
# S7: Case Package
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "Case Package 契约", "标注工具与平台之间的离线文件交换标准")

# Left: directory tree
tree = add_rect(s, Inches(0.5), Inches(1.6), Inches(5.0), Inches(3.0), fill_color=DARK)
tree_lines = [
    ("case_package/", True, 14, BLUE),
    ("├── manifest.json（sha256 / shape / spacing）", False, 12, GRAY),
    ("├── images/image.nii.gz", False, 12, GRAY),
    ("├── labels/", False, 12, WHITE),
    ("│   ├── draft_label.nii.gz", False, 12, GRAY),
    ("│   ├── verified_label.nii.gz", False, 12, GREEN),
    ("│   └── masks/{liver, kidney...}.nii.gz", False, 12, GRAY),
    ("├── config/{anatomy_vocabulary, review_label_map}.yaml", False, 12, GRAY),
    ("├── reports/{geometry_check, review_report}.json", False, 12, GRAY),
    ("└── provenance/{source_labels, tool_export}.json", False, 12, GRAY),
]
add_multiline(s, Inches(0.7), Inches(1.7), Inches(4.5), Inches(2.8), tree_lines)

# Right: validation rules
rows = [
    [("校验规则", True, WHITE, DARK), ("级别", True, WHITE, DARK), ("处理方式", True, WHITE, DARK)],
    ["图像 sha256 不一致", "Error", "拒绝导入"],
    ["标签 shape 不一致", "Error", "不可修复，拒绝"],
    ["spacing/affine 不一致", "Warning", "check_geometry.py 自动修复"],
    ["label id 不合法", "Error", "拒绝"],
    ["只有 draft 无 verified", "Warning", "标记，不阻塞"],
]
add_table(s, Inches(6.0), Inches(1.6), [Inches(3.0), Inches(1.5), Inches(2.5)], rows)

# Bottom
add_text(s, Inches(0.5), Inches(5.2), Inches(12), Inches(0.5),
         "一个病例一个包，自包含可搬运。nnUNet 不直接读 Case Package → 经 Data Registry + Snapshot 导出",
         font_size=11, color=DGRAY)

# ═══════════════════════════════════════
# S8: nnUNet Pipeline
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "训练管线 — nnUNet 五阶段")

actions = [
    ("Action1\n数据转换", "标注数据→nnUNet格式\n重采样+合并+划分", BLUE),
    ("Action2\n预处理", "指纹提取+实验规划\n支持手动覆盖参数", BLUE),
    ("Action3\n训练", "GPU 管理+5折交叉\nDDP 多卡支持", GREEN),
    ("Action4\n推理", "预插值加速\n多模型共享分辨率", AMBER),
    ("Action5\n评估", "Dice+Surface Dice\n多格式报告+聚合", AMBER),
]
aw = Inches(2.25)
ah = Inches(1.5)
for i, (title, desc, col) in enumerate(actions):
    x = Inches(0.4) + i * (aw + Inches(0.25))
    y = Inches(1.6)
    add_rect(s, x, y, aw, ah, fill_color=WHITE, line_color=col, line_width=Pt(1.5))
    add_text(s, x + Inches(0.1), y + Inches(0.1), aw - Inches(0.2), Inches(0.6),
             title, font_size=12, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.75), aw - Inches(0.2), Inches(0.65),
             desc, font_size=9, color=DGRAY, alignment=PP_ALIGN.CENTER)
    if i < len(actions) - 1:
        add_text(s, x + aw - Inches(0.05), y + Inches(0.55), Inches(0.35), Inches(0.3),
                 "→", font_size=20, color=BLUE, alignment=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
         "Config.toml + ModelMap.toml → AutoSegmentationFramework（编排器）",
         font_size=14, bold=True, color=BLACK)
add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(1.5),
         "ModelMap.toml 已是 TaskLabelMap 雏形 — 每模型独立从 1 编号。支持扁平格式（细粒度）和分组格式（粗分割）\n多 GPU 调度：不同 GPU 并行 + 同 GPU 串行 → 自动生成 gpu{0,1,2,3}.sh 脚本",
         font_size=11, color=DGRAY)

# ═══════════════════════════════════════
# S9: Adapter Architecture
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "Adapter 架构 — 可扩展性核心")

add_text(s, Inches(4.2), Inches(1.5), Inches(4.9), Inches(0.5),
         "Dataset Snapshot（冻结的数据视图）",
         font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

adapters_data = [
    ("nnUNet Adapter", "全监督训练\n生产级 · 已完成", GREEN),
    ("FewShot Adapter", "预训练 + 微调\n架构已定 · 待验证", BLUE),
    ("MONAI Adapter", "Transformer 模型\n预留", AMBER),
    ("其他 Adapter", "SAM / 新框架\n可扩展", GRAY),
]
aw = Inches(2.8)
ah = Inches(1.3)
for i, (title, desc, col) in enumerate(adapters_data):
    x = Inches(0.5) + i * (aw + Inches(0.35))
    y = Inches(2.4)
    add_rect(s, x, y, aw, ah, fill_color=WHITE, line_color=col, line_width=Pt(1.5))
    add_text(s, x + Inches(0.1), y + Inches(0.1), aw - Inches(0.2), Inches(0.5),
             title, font_size=13, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.65), aw - Inches(0.2), Inches(0.55),
             desc, font_size=10, color=DGRAY, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), fill_color=WHITE)
lines = [
    ("不变层（数据契约）", True, 16, BLUE),
    ("Case / Image Artifact / Label Artifact / Dataset Snapshot / label_policy", False, 13, BLACK),
    ("", False, 6, BLACK),
    ("可变层（Adapter 封装）", True, 16, GREEN),
    ("nnUNet / MONAI / FewShot / SAM / ...", False, 13, BLACK),
    ("", False, 6, BLACK),
    ("新框架 = 新 Adapter，不动数据契约。USB 协议：设备可换，接口不变。", False, 13, DGRAY),
]
add_multiline(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.1), lines)

# ═══════════════════════════════════════
# S10: FewShot Learning
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "少样本学习 — FewShot Adapter")

rows = [
    [("架构定位", True, WHITE), ("training 域新 Adapter，与 nnUNet Adapter 平级。先通过实验协议验证 → 升级为正式 Adapter", False, WHITE, DARK)],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(2.5), Inches(10.0)], rows)

rows2 = [
    [("实验协议", True, WHITE), ("说明", True, WHITE)],
    ["1. 选定器官", "从 Registry 选出 verified_label 病例，患者级别冻结划分"],
    ["2. N-shot Snapshot", "构建 N=1/3/5/10/20 的快照"],
    ["3. 三对照组", "A. 全监督上界 / B. 同数据量基线 / C. 预训练+微调实验组"],
    ["4. 冻结评估", "冻结测试集上计算 Dice, Surface Dice, Hausdorff"],
    ["5. 登记", "Model Registry 记录 N、器官、标签状态、数据版本"],
]
add_table(s, Inches(0.5), Inches(2.5), [Inches(2.5), Inches(10.0)], rows2)

rows3 = [
    [("准入标准", True, WHITE, DARK), ("阈值", True, WHITE, DARK)],
    ["独立复现", "每器官 ≥ 3 次"],
    ["Dice vs 全监督", "≥ 全监督的 90%"],
    ["跨扫描协议差距", "≤ 0.05 Dice"],
    ["失败率（Dice < 0.3）", "< 5%"],
]
add_table(s, Inches(0.5), Inches(5.5), [Inches(4.0), Inches(4.0)], rows3)

# ═══════════════════════════════════════
# S11: Mimics Integration
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "Mimics 集成方案")

add_text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.4),
         "流程：平台 split/merge → Case Package → Mimics 导入 → 人工修正 → Mimics 导出 → 校验",
         font_size=13, bold=True, color=BLACK)

rows = [
    [("API", True, WHITE), ("用途", True, WHITE), ("来源", True, WHITE)],
    ["mimics.data.masks.find(name=...)", "查找指定名称 Mask", "社区代码"],
    ["mask.get_voxel_buffer()", "获取 Mask 体素（numpy 数组）", "社区代码"],
    ["mask.set_voxel_buffer(arr)", "导入外部 numpy 数组到 Mask", "官方员工确认 (2023)"],
    ["NIfTI 导入导出（2025 GUI）", "官方原生功能，替代手工写 affine", "产品更新页"],
    ["Help → Scripting Guide", "完整 API 文档，内置在 Mimics 中", "随软件安装"],
]
add_table(s, Inches(0.5), Inches(2.1), [Inches(4.0), Inches(5.0), Inches(3.5)], rows)

rows2 = [
    [("风险", True, WHITE), ("影响", True, WHITE), ("应对", True, WHITE)],
    ["set_voxel_buffer 空间对齐", "标签与 CT 错位", "check_geometry.py 检测 + 修复"],
    ["Scripting API 不足以覆盖所有步骤", "自动化降低", "GUI 手动 + 平台脚本补偿"],
    ["Mimics 完全不可用", "阻塞标注流程", "切换 3D Slicer / ITK-SNAP"],
]
add_table(s, Inches(0.5), Inches(5.5), [Inches(3.5), Inches(4.0), Inches(5.0)], rows2)

# ═══════════════════════════════════════
# S12: Implementation Phases
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "分阶段实施路线")

rows = [
    [("阶段", True, WHITE), ("目标", True, WHITE), ("状态", True, WHITE)],
    [("A. 文件包闭环", True, BLACK, LGRAY), "一个病例从候选标签到人工保存，再进入训练快照", ("👈 当前", True, GREEN)],
    ["B. 注册中心 + 快照", "Data Registry + Dataset Snapshot，多任务复用，训练可复现", "设计已定"],
    ["C. Adapter 稳定", "nnUNet Adapter 跑通，预留 MONAI / FewShot 接口", "设计已定"],
    ["D. 离线批量推理", "模型版本 + 批量任务 + candidate_label 回流", "后期"],
    ["E. 统一调度", "Web UI、任务队列、权限、审计 → 数据契约稳定后", "后期"],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(3.0), Inches(7.3), Inches(2.2)], rows)

# ═══════════════════════════════════════
# S13: Key Decisions
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "已确认的关键决策")

rows = [
    [("决策", True, WHITE), ("决定", True, WHITE)],
    ["伪标签准入", "默认允许，取决于器官/任务/模型。默认允许 + 特定排除"],
    ["器官范围", "v500 全部模型（CT1-16、MR1-8），全身覆盖"],
    ["全身模型方案", "保持多模型组合（已验证），后期再实验统一模型"],
    ["Mimics POC", "先调研准备（确认 Scripting Guide API），再启动 POC"],
    ["少样本学习", "training 域新 Adapter，与 nnUNet 平级，需过实验协议"],
    ["label_policy", "candidate_label 可通过 allow_status 直接纳入训练"],
    ["Data Registry", "先不实现但必记，闭环跑通后补"],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(2.5), Inches(10.0)], rows)

# ═══════════════════════════════════════
# S14: Checklist
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, LGRAY)
slide_title_bar(s, "实现前待确认清单")

rows = [
    [("阻塞项", True, WHITE, RED), ("状态", True, WHITE, RED)],
    ["A1. Mimics 版本 + 许可证类型", "未确认"],
    ["A2. Scripting Guide 中 mask 相关 API 清单", "需在 Mimics Help 中查看"],
    ["A3. Mimics 能否运行 Python 脚本", "取决于许可证"],
    ["A4. 第一批 3-5 个病例", "数据不是瓶颈"],
    ["A5. 训练服务器环境（GPU、路径）", "未确认"],
]
add_table(s, Inches(0.5), Inches(1.6), [Inches(6.0), Inches(6.3)], rows)

rows2 = [
    [("非阻塞项（可立即实现）", True, WHITE, GREEN), ("用途", True, WHITE, GREEN)],
    ["B1. split_multilabel_to_masks.py", "多标签 NIfTI → 逐器官二值 mask"],
    ["B2. merge_masks_to_multilabel.py", "逐器官 mask → 多标签 NIfTI"],
    ["B3. check_geometry.py", "shape/spacing/affine 校验 + 自动修复"],
    ["B4. package_case.py", "生成 Case Package 目录"],
    ["B5. import_case_package.py（Mimics Adapter）", "读 masks → numpy → set_voxel_buffer"],
    ["B6. export_review_package.py（Mimics Adapter）", "get_voxel_buffer → numpy → 保存 masks"],
]
add_table(s, Inches(0.5), Inches(4.0), [Inches(6.0), Inches(6.3)], rows2)

# ═══════════════════════════════════════
# S15: Time Estimate
# ═══════════════════════════════════════
s = add_blank_slide()
add_bg(s, DARK)
add_rect(s, 0, 0, W, Inches(0.06), fill_color=BLUE)
add_text(s, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "时间估算",
         font_size=28, bold=True, color=WHITE)

# Three scenario cards
scenarios = [
    ("顺利", "3-4 周", "A1-A3 秒确认\nAPI 全有，几何无问题", GREEN),
    ("正常", "5-6 周", "1-2 个 Mimics 弯路\n1-2 轮几何调试", AMBER),
    ("Mimics 受限", "+2-3 周", "降级 GUI 手动\n或切 3D Slicer", RED),
]
sw = Inches(3.8)
for i, (title, time_val, desc, col) in enumerate(scenarios):
    x = Inches(0.5) + i * (sw + Inches(0.4))
    y = Inches(1.3)
    add_rect(s, x, y, sw, Inches(2.3), fill_color=RGBColor(0x1E, 0x29, 0x3B))
    add_text(s, x + Inches(0.2), y + Inches(0.15), sw - Inches(0.4), Inches(0.4),
             title, font_size=14, bold=True, color=col, alignment=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(0.6), sw - Inches(0.4), Inches(0.8),
             time_val, font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(1.0), y + Inches(1.45), Inches(1.8), Inches(0.03), fill_color=col)
    add_text(s, x + Inches(0.2), y + Inches(1.6), sw - Inches(0.4), Inches(0.6),
             desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Breakdown table
rows = [
    [("模块", True, WHITE), ("工时", True, WHITE), ("信心", True, WHITE)],
    ["平台脚本 B1-B4", "4-5 天", "高"],
    ["Mimics Adapter B5-B6", "乐观 2 天 / 悲观 2 周", "中低（依赖 A1-A3）"],
    ["闭环测试 + 调试", "1-2 周", "—"],
    ["3-5 病例验证", "1 周", "—"],
]
add_table(s, Inches(0.5), Inches(4.0), [Inches(5.0), Inches(4.0), Inches(3.3)], rows)

add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
         "最大减速带不是代码量，是 Mimics 确认和几何对齐调试",
         font_size=11, color=GRAY)

# ── Save ──
output_path = "/Users/ruanshijian/SegmentationPlatform/SegmentationPlatform_Architecture.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
